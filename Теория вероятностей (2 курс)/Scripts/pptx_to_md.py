#!/usr/bin/env python3
"""Конвертация PPTX-лекций в Markdown (без вступительного слайда с контактами)."""

import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation

BASE = Path(__file__).resolve().parents[1]
SRC = BASE / "Презентации"
OUT = BASE / "конспекты"
M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"

JUNK_LINES = {
    "в", "а", "b", "Ω", "ω", "произошло", "не произошло", "аб", "ав", "в",
    "область построения диаграммы",
}

KNOWN_FORMULAS: list[tuple[str, str]] = [
    (r"формул[аеы].*бернулли|по формуле бернулли", r"$$P_n(m) = C_n^m\, p^m q^{n-m}, \quad q = 1-p$$"),
    (r"условн.*вероятност", r"$$P(B \mid A) = \frac{P(AB)}{P(A)}, \quad P(A) > 0$$"),
    (r"формул[ае].*умножен", r"$$P(AB) = P(A)\,P(B \mid A) = P(B)\,P(A \mid B)$$"),
    (r"полной вероятност", r"$$P(A) = \sum_{i=1}^{n} P(H_i)\,P(A \mid H_i)$$"),
    (r"байес", r"$$P(H_k \mid A) = \frac{P(H_k)\,P(A \mid H_k)}{\sum_i P(H_i)\,P(A \mid H_i)}$$"),
    (r"неравенств[оа].*марков", r"$$P(X \geq a) \leq \frac{M[X]}{a}, \quad a > 0$$"),
    (r"неравенств[оа].*чебыш", r"$$P(|X - M[X]| \geq \varepsilon) \leq \frac{D[X]}{\varepsilon^2}$$"),
    (r"математическ.*ожидан", r"$$M[X] = \sum_i x_i p_i \quad \text{(дискр.)}; \qquad M[X] = \int_{-\infty}^{\infty} x f(x)\,dx \quad \text{(непр.)}$$"),
    (r"дисперс", r"$$D[X] = M[X^2] - (M[X])^2, \qquad D[aX+b] = a^2 D[X]$$"),
    (r"выборочн.*средн", r"$$\bar{x} = \frac{1}{n}\sum_{i=1}^{k} x_i n_i$$"),
    (r"ряд распределен", r"$$\sum_i p_i = 1$$"),
    (r"нормальн.*закон|нормальн.*распределен", r"$$f(x) = \frac{1}{\sigma\sqrt{2\pi}} \exp\!\left(-\frac{(x-\mu)^2}{2\sigma^2}\right)$$"),
    (r"равномерн.*закон|равномерн.*распределен", r"$$f(x) = \frac{1}{b-a}, \quad x \in [a,b]$$"),
    (r"плотност.*вероятност", r"$$F(x) = \int_{-\infty}^{x} f(t)\,dt, \quad f(x) = F'(x)$$"),
    (r"ковариац", r"$$\mathrm{cov}(X,Y) = M[XY] - M[X]M[Y]$$"),
    (r"коэффициент корреляц", r"$$r_{XY} = \frac{\mathrm{cov}(X,Y)}{\sigma_X \sigma_Y}$$"),
    (r"закон больших чисел|теорема чебыш.*больш", r"$$\lim_{n\to\infty} P\!\left(\left|\frac{1}{n}\sum_{i=1}^n X_i - M[X]\right| \geq \varepsilon\right) = 0$$"),
    (r"теорема бернулли", r"$$\lim_{n\to\infty} P\!\left(\left|\frac{m}{n} - p\right| \geq \varepsilon\right) = 0$$"),
    (r"центральн.*предельн|теорема ляпунов", r"$$\frac{\sum X_i - n\mu}{\sigma\sqrt{n}} \xrightarrow{d} N(0,1)$$"),
]


def lecture_num(name: str) -> int:
    m = re.search(r"Лекция\s+(\d+)", name)
    return int(m.group(1)) if m else 0


def is_intro_slide(text: str) -> bool:
    t = text.lower()
    markers = ["преподаватель", "e-mail", "аттестация", "nntu@yandex"]
    return sum(m in t for m in markers) >= 2


def is_closing_slide(text: str) -> bool:
    t = text.lower()
    return "спасибо за внимание" in t or "на сегодня все" in t


def clean_line(line: str) -> str | None:
    line = line.strip()
    if not line:
        return None
    if re.fullmatch(r"\d{1,3}", line):
        return None
    if line.lower() in JUNK_LINES:
        return None
    if re.fullmatch(r"[А-ЯA-Zа-яA-Za-z]{1,2}", line):
        return None
    line = line.replace("Ѳ", "θ").replace("ϴ", "Θ").replace("–", "—")
    return line


def clean_text(text: str) -> str:
    lines = [cl for l in text.splitlines() if (cl := clean_line(l)) is not None]
    return "\n".join(lines)


def shape_texts(slide) -> list[str]:
    chunks = []
    for shape in slide.shapes:
        if shape.has_table:
            rows = []
            for row in shape.table.rows:
                rows.append("| " + " | ".join(c.text.strip() for c in row.cells) + " |")
            if rows:
                sep = "| " + " | ".join("---" for _ in shape.table.rows[0].cells) + " |"
                chunks.append("\n".join([rows[0], sep] + rows[1:]))
        elif hasattr(shape, "text") and shape.text.strip():
            chunks.append(shape.text.strip())
    return chunks


def omml_element_text(el) -> str:
    tag = el.tag.split("}")[-1] if "}" in el.tag else el.tag
    if tag == "t":
        return (el.text or "") + (el.tail or "")
    if tag == "f":
        num = den = ""
        for ch in el:
            ct = ch.tag.split("}")[-1]
            if ct == "num":
                num = "".join(omml_element_text(x) for x in ch.iter())
            elif ct == "den":
                den = "".join(omml_element_text(x) for x in ch.iter())
        return f"\\frac{{{num}}}{{{den}}}"
    if tag == "sSup":
        parts = ["".join(omml_element_text(x) for x in ch.iter()) for ch in el]
        if len(parts) >= 2:
            return f"{parts[0]}^{{{parts[1]}}}"
    if tag == "sSub":
        parts = ["".join(omml_element_text(x) for x in ch.iter()) for ch in el]
        if len(parts) >= 2:
            return f"{parts[0]}_{{{parts[1]}}}"
    return "".join(omml_element_text(c) for c in el)


def extract_omml(pptx_path: Path, slide_idx: int) -> list[str]:
    with zipfile.ZipFile(pptx_path) as z:
        xml = z.read(f"ppt/slides/slide{slide_idx + 1}.xml")
    root = ET.fromstring(xml)
    out = []
    for om in root.iter(f"{{{M_NS}}}oMath"):
        s = omml_element_text(om).strip()
        if s:
            out.append(f"$$ {s} $$")
    return out


def known_formulas(text: str) -> list[str]:
    low = text.lower()
    if len(low) < 40:
        return []
    found = []
    for pattern, formula in KNOWN_FORMULAS:
        if re.search(pattern, low) and formula not in found:
            found.append(formula)
    if re.search(r"p\(а\|в\)|p\(a\|b\)", low):
        f = r"$$P(A \mid B) = \frac{P(AB)}{P(B)}, \quad P(B) > 0$$"
        if f not in found:
            found.append(f)
    if re.search(r"следовательно\s*,?\s*$", low) or "полная группа" in low:
        f = r"$$\sum_i p_i = 1$$"
        if f not in found:
            found.append(f)
    return found


def slide_has_many_images(pptx_path: Path, slide_idx: int) -> bool:
    with zipfile.ZipFile(pptx_path) as z:
        xml = z.read(f"ppt/slides/slide{slide_idx + 1}.xml").decode("utf-8", errors="ignore")
    return xml.count("a:blip") > 2


def slide_to_md(slide, pptx_path: Path, slide_idx: int) -> tuple[str, str]:
    text = clean_text("\n\n".join(shape_texts(slide)))
    if not text:
        omml = extract_omml(pptx_path, slide_idx)
        return ("\n\n".join(omml), "content") if omml else ("", "empty")

    if is_intro_slide(text) or is_closing_slide(text):
        return "", "skip"

    formulas = extract_omml(pptx_path, slide_idx) + known_formulas(text)
    lines = text.split("\n")
    title, body = lines[0], lines[1:]

    is_heading = (
        len(lines) <= 4
        and len(title) < 100
        and not title.endswith(":")
        and (not body or all(len(b) < 90 for b in body))
    )

    if is_heading:
        parts = [f"## {title}"]
        if body:
            parts.append("\n".join(body))
    elif title.endswith(":"):
        parts = [text]
    else:
        parts = [f"## {title}"] + body if title else [text]

    block = "\n\n".join(parts)
    for f in formulas:
        if f not in block:
            block += f"\n\n{f}"

    if slide_has_many_images(pptx_path, slide_idx) and not extract_omml(pptx_path, slide_idx):
        if re.search(r"формул|решение|равн|неравенств|теорем", text.lower()):
            block += "\n\n> *Часть содержимого — графические формулы в PPTX.*"

    return block, "content"


def extract_topic(blocks: list[str], lec_num: int) -> str:
    if blocks:
        m = re.match(r"##\s+(.+)", blocks[0])
        if m:
            t = m.group(1).strip()
            if "лекция" not in t.lower() and len(t) < 80:
                return t
    return f"Тема лекции {lec_num}"


def convert_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    lec = lecture_num(path.name)
    blocks = []
    for i in range(len(prs.slides)):
        block, kind = slide_to_md(prs.slides[i], path, i)
        if kind == "content" and block.strip():
            blocks.append(block)

    topic = extract_topic(blocks, lec)
    if blocks and blocks[0].startswith("## "):
        first_line, rest = (
            (blocks[0].split("\n", 1)[0][3:].strip(), blocks[0].split("\n", 1)[1].strip())
            if "\n" in blocks[0]
            else (blocks[0][3:].strip(), "")
        )
        if first_line.lower() == topic.lower():
            blocks[0] = rest
    blocks = [b for b in blocks if b.strip() and not re.fullmatch(r"(\$\$.*\$\$\s*)+", b, flags=re.DOTALL)]

    body = "\n\n---\n\n".join(blocks)
    return (
        f"# Лекция {lec}. {topic}\n\n{body}\n\n"
        f"---\n\n*Источник: `{path.name}` · "
        f"формулы с картинок дополнены по контексту, где текст не извлёкся.*\n"
    )


def main():
    OUT.mkdir(exist_ok=True)
    for pptx in sorted(SRC.glob("*.pptx")):
        n = lecture_num(pptx.name)
        (OUT / f"лекция_{n:02d}.md").write_text(convert_pptx(pptx), encoding="utf-8")
        print(f"OK: лекция_{n:02d}.md")


if __name__ == "__main__":
    main()
